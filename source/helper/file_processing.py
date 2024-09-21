# helper/file_processing.py

import os
from langchain.document_loaders import PyPDFLoader  # For loading PDF files
from langchain_experimental.text_splitter import SemanticChunker  # For splitting text
from langchain.embeddings.openai import OpenAIEmbeddings  # For generating embeddings
from langchain_pinecone import PineconeVectorStore  # For integrating LangChain with Pinecone
from helper.config import get_settings
import pinecone
from pptx import Presentation  # For processing PowerPoint files

# Initialize OpenAI Embeddings
def initialize_embeddings():
    settings = get_settings()
    embeddings = OpenAIEmbeddings(model="text-embedding-ada-002", openai_api_key=settings.OPENAI_API_KEY)
    return embeddings

# Initialize Pinecone Index
def initialize_pinecone():
    settings = get_settings()
    pinecone.init(api_key=settings.PINECONE_API_KEY, environment=settings.PINECONE_ENV)
    return pinecone

# Function to read TXT files
def read_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    return content

# Function to process PowerPoint (PPTX) files
def process_pptx(file_path):
    prs = Presentation(file_path)
    content = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text)

    return "\n".join(content)

# Main function to process file, generate embeddings, and add them to the index
def process_file(file_path, file_type):
    settings = get_settings()
    embeddings = initialize_embeddings()
    documents = []

    # Load content based on file type
    if file_type == 'pdf':
        loader = PyPDFLoader(file_path)
        documents = loader.load()
    elif file_type == 'txt':
        text_content = read_txt(file_path)
        documents = [{"text": text_content}]
    elif file_type == 'pptx':
        text_content = process_pptx(file_path)
        documents = [{"text": text_content}]
    else:
        raise ValueError("Unsupported file type.")

    # Split documents into smaller chunks for embedding
    text_splitter = SemanticChunker(embeddings, breakpoint_threshold_type="percentile")
    split_docs = text_splitter.split_documents(documents)

    # Connect or create the Pinecone index
    pinecone_instance = initialize_pinecone()
    index_name = settings.PINECONE_INDEX_NAME
    if index_name not in pinecone_instance.list_indexes():
        pinecone_instance.create_index(name=index_name, dimension=settings.PINECONE_DIMENSION, metric=settings.PINECONE_METRIC)

    # Populate the Pinecone index with document embeddings
    vectorstore = PineconeVectorStore.from_documents(split_docs, embeddings, index_name=index_name)

    return vectorstore
